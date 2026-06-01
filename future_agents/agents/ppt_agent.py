"""PPT Agent — creates beautiful, professionally structured PowerPoint presentations.

Uses python-pptx with light, vibrant color themes. Auto-generates full structure:
Title → Index → Introduction → Topic slides → Architecture → Summary → Q&A.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any

from future_agents.core.base_agent import BaseAgent, TaskContext, TaskResult
from future_agents.models.feedback import ExecutionOutcome

logger = logging.getLogger(__name__)

# ── Light color palette ──────────────────────────────────────────────────────
PALETTE = {
    "sky":      {"bg": (0xE3, 0xF2, 0xFD), "accent": (0x19, 0x76, 0xD2), "text": (0x0D, 0x47, 0xA1)},
    "mint":     {"bg": (0xE8, 0xF5, 0xE9), "accent": (0x2E, 0x7D, 0x32), "text": (0x1B, 0x5E, 0x20)},
    "lavender": {"bg": (0xF3, 0xE5, 0xF5), "accent": (0x7B, 0x1F, 0xA2), "text": (0x4A, 0x14, 0x8C)},
    "sunrise":  {"bg": (0xFF, 0xF8, 0xE1), "accent": (0xF5, 0x7F, 0x17), "text": (0xE6, 0x5C, 0x00)},
    "coral":    {"bg": (0xFF, 0xEB, 0xEE), "accent": (0xC6, 0x28, 0x28), "text": (0xB7, 0x1C, 0x1C)},
    "teal":     {"bg": (0xE0, 0xF2, 0xF1), "accent": (0x00, 0x69, 0x6C), "text": (0x00, 0x4D, 0x40)},
}

SLIDE_THEMES = [
    "sky", "mint", "lavender", "sunrise", "coral", "teal",
    "sky", "mint", "lavender", "sunrise",
]


def _rgb(r: int, g: int, b: int):
    from pptx.util import Pt  # noqa — imported here to keep top-level import clean
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _pt(n: float):
    from pptx.util import Pt
    return Pt(n)


def _emu(cm: float):
    from pptx.util import Cm
    return Cm(cm)


class PPTAgent(BaseAgent):
    """Creates stunning PowerPoint presentations with structured, light-themed slides.

    Capabilities:
    - ppt.create  — full deck from topic + outline
    - ppt.export  — get the file bytes for download
    """

    def __init__(self, output_dir: str = "output/ppt", **kwargs: Any) -> None:
        super().__init__(**kwargs)
        self._output_dir = Path(output_dir)
        self._output_dir.mkdir(parents=True, exist_ok=True)

    @property
    def agent_type(self) -> str:
        return "ppt"

    @property
    def capabilities(self) -> list[str]:
        return ["ppt.create", "ppt.export"]

    async def _execute(self, context: TaskContext) -> TaskResult:
        handlers = {
            "ppt.create": self._create_presentation,
            "ppt.export": self._export_presentation,
        }
        handler = handlers.get(context.intent)
        if not handler:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=[f"Unknown intent: {context.intent}"],
            )
        return await handler(context)

    # ── Public API ───────────────────────────────────────────────────────────

    async def _create_presentation(self, context: TaskContext) -> TaskResult:
        """Build a full presentation deck.

        Parameters:
            title        (str): Presentation title
            subtitle     (str): Optional subtitle / speaker name
            topics       (list[dict]): Each has {"title": str, "bullets": [str], "notes": str}
            color_theme  (str): One of sky|mint|lavender|sunrise|coral|teal (default: sky)
            output_name  (str): Output filename without extension
            author       (str): Author name for footer
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Cm
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=["python-pptx not installed. Run: pip install python-pptx"],
            )

        p = context.parameters
        title = p.get("title", "Untitled Presentation")
        subtitle = p.get("subtitle", "")
        topics = p.get("topics", [])
        theme_name = p.get("color_theme", "sky")
        output_name = p.get("output_name", title.lower().replace(" ", "_")[:40])
        author = p.get("author", "Document Agent")

        theme = PALETTE.get(theme_name, PALETTE["sky"])
        bg_color = RGBColor(*theme["bg"])
        accent = RGBColor(*theme["accent"])
        text_color = RGBColor(*theme["text"])

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank

        def add_bg(slide, color: RGBColor):
            from pptx.util import Inches
            bg = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                0, 0, prs.slide_width, prs.slide_height,
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()
            bg.shadow.inherit = False
            return bg

        def add_accent_bar(slide, color: RGBColor, height_cm: float = 0.8):
            bar = slide.shapes.add_shape(
                1, 0, 0, prs.slide_width, _emu(height_cm),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            return bar

        def text_box(slide, left, top, width, height, text, size, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
            from pptx.util import Inches, Pt
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color or text_color
            return txb

        # ── SLIDE 1: Title ──────────────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, bg_color)

        # Decorative top bar
        top_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, _emu(1.5))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = accent
        top_bar.line.fill.background()

        # Bottom bar
        bot_bar = slide.shapes.add_shape(1, 0, prs.slide_height - _emu(0.6), prs.slide_width, _emu(0.6))
        bot_bar.fill.solid()
        bot_bar.fill.fore_color.rgb = accent
        bot_bar.line.fill.background()

        # Title text
        text_box(
            slide,
            Inches(1), Inches(1.8), Inches(11.33), Inches(1.8),
            title, 44, bold=True, color=text_color, align=PP_ALIGN.CENTER,
        )

        if subtitle:
            text_box(
                slide,
                Inches(2), Inches(3.7), Inches(9.33), Inches(0.8),
                subtitle, 22, italic=True, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER,
            )

        # Author / date
        text_box(
            slide,
            Inches(1), Inches(5.0), Inches(11.33), Inches(0.5),
            f"Presented by {author}", 14, color=RGBColor(0x77, 0x77, 0x77), align=PP_ALIGN.CENTER,
        )

        # ── SLIDE 2: Table of Contents ──────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, bg_color)
        add_accent_bar(slide, accent)

        text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 "Table of Contents", 22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        toc_items = ["Introduction"] + [t.get("title", f"Topic {i+1}") for i, t in enumerate(topics)] + [
            "Architecture Overview", "Summary", "Q & A",
        ]
        col_break = (len(toc_items) + 1) // 2

        for idx, item in enumerate(toc_items):
            col = idx // col_break
            row = idx % col_break
            lx = Inches(0.8 + col * 6.2)
            ly = Inches(1.2 + row * 0.65)
            # Bullet circle
            circ = slide.shapes.add_shape(9, lx, ly + _emu(0.15), _emu(0.5), _emu(0.5))
            circ.fill.solid()
            circ.fill.fore_color.rgb = accent
            circ.line.fill.background()

            # Number inside circle
            tf = circ.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = str(idx + 1)
            run.font.size = _pt(8)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            text_box(slide, lx + _emu(0.7), ly, Inches(5.5), _emu(0.65),
                     item, 14, color=text_color)

        # ── SLIDE 3: Introduction ───────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, bg_color)
        add_accent_bar(slide, accent)
        text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 "Introduction", 22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        intro_text = p.get("introduction", f"This presentation covers key aspects of: {title}.")
        text_box(slide, Inches(0.8), Inches(1.2), Inches(11.6), Inches(5.5),
                 intro_text, 18, color=text_color)

        # ── SLIDES 4..N: Topics ─────────────────────────────────────
        for i, topic in enumerate(topics):
            t_theme = PALETTE[SLIDE_THEMES[i % len(SLIDE_THEMES)]]
            t_bg = RGBColor(*t_theme["bg"])
            t_accent = RGBColor(*t_theme["accent"])
            t_text = RGBColor(*t_theme["text"])

            slide = prs.slides.add_slide(blank_layout)
            add_bg(slide, t_bg)

            # Side accent stripe
            stripe = slide.shapes.add_shape(1, 0, 0, _emu(0.6), prs.slide_height)
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = t_accent
            stripe.line.fill.background()

            # Topic number badge
            badge = slide.shapes.add_shape(9, _emu(0.2), _emu(0.3), _emu(1.2), _emu(1.2))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            badge.line.fill.background()
            tf = badge.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = str(i + 1)
            run.font.size = _pt(16)
            run.font.bold = True
            run.font.color.rgb = t_accent

            # Title
            text_box(slide, Inches(1.0), Inches(0.3), Inches(11.5), Inches(0.9),
                     topic.get("title", ""), 28, bold=True, color=t_text)

            # Divider line
            divider = slide.shapes.add_shape(1, Inches(1.0), Inches(1.35), Inches(11.3), _emu(0.06))
            divider.fill.solid()
            divider.fill.fore_color.rgb = t_accent
            divider.line.fill.background()

            # Bullets
            bullets = topic.get("bullets", [])
            for j, bullet in enumerate(bullets[:7]):
                by = Inches(1.55) + j * Inches(0.75)
                dot = slide.shapes.add_shape(9, Inches(1.0), by + _emu(0.25), _emu(0.4), _emu(0.4))
                dot.fill.solid()
                dot.fill.fore_color.rgb = t_accent
                dot.line.fill.background()
                text_box(slide, Inches(1.6), by, Inches(11.0), Inches(0.7),
                         bullet, 15, color=RGBColor(0x33, 0x33, 0x33))

            # Speaker notes
            notes = topic.get("notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        # ── Architecture Slide ──────────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, RGBColor(0xF0, 0xF4, 0xFF))
        add_accent_bar(slide, RGBColor(0x19, 0x76, 0xD2))
        text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 "Architecture Overview", 22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        arch_boxes = p.get("architecture", [
            {"label": "Input Layer", "color": (0x42, 0xA5, 0xF5)},
            {"label": "Processing Engine", "color": (0x66, 0xBB, 0x6A)},
            {"label": "Output Layer", "color": (0xFFA7, 0x26, 0x00)},
        ])

        box_w = Inches(2.8)
        box_h = Inches(1.2)
        gap = Inches(0.6)
        total_w = len(arch_boxes) * box_w + (len(arch_boxes) - 1) * gap
        start_x = (prs.slide_width - total_w) / 2
        box_y = Inches(2.5)

        for idx_a, abox in enumerate(arch_boxes):
            bx = start_x + idx_a * (box_w + gap)
            color_rgb = RGBColor(*abox.get("color", (0x42, 0xA5, 0xF5)))
            rect = slide.shapes.add_shape(1, bx, box_y, box_w, box_h)
            rect.fill.solid()
            rect.fill.fore_color.rgb = color_rgb
            rect.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf = rect.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = abox.get("label", f"Component {idx_a + 1}")
            run.font.size = _pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Arrow between boxes
            if idx_a < len(arch_boxes) - 1:
                arr_x = bx + box_w
                arr_y = box_y + box_h / 2 - _emu(0.15)
                arr = slide.shapes.add_shape(13, arr_x, arr_y, gap, _emu(0.3))
                arr.fill.solid()
                arr.fill.fore_color.rgb = RGBColor(0x90, 0xA4, 0xAE)
                arr.line.fill.background()

        arch_desc = p.get("architecture_description", "High-level system component flow.")
        text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1.0),
                 arch_desc, 14, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

        # ── Summary Slide ───────────────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, bg_color)
        add_accent_bar(slide, accent)
        text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 "Summary", 22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

        summary_points = p.get("summary", [
            f"Key insights from: {title}",
            "Review all topic highlights",
            "Actionable next steps defined",
        ])
        for idx_s, sp in enumerate(summary_points[:6]):
            sy = Inches(1.2) + idx_s * Inches(0.85)
            chk = slide.shapes.add_shape(1, Inches(0.7), sy + _emu(0.1), _emu(0.55), _emu(0.55))
            chk.fill.solid()
            chk.fill.fore_color.rgb = accent
            chk.line.fill.background()
            tf = chk.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = "✓"
            run.font.size = _pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            text_box(slide, Inches(1.4), sy, Inches(11.2), Inches(0.8),
                     sp, 16, color=text_color)

        # ── Q&A Slide ───────────────────────────────────────────────
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, bg_color)

        # Large decorative Q
        big_q = slide.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        tf = big_q.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = "?"
        run.font.size = _pt(300)
        run.font.bold = True
        run.font.color.rgb = RGBColor(*[min(c + 30, 255) for c in theme["bg"]])

        text_box(slide, Inches(0), Inches(2.5), prs.slide_width, Inches(1.5),
                 "Questions & Answers", 40, bold=True, color=accent, align=PP_ALIGN.CENTER)
        text_box(slide, Inches(1), Inches(4.2), Inches(11.33), Inches(0.6),
                 "Thank you for your attention!", 18,
                 color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

        # ── Save ────────────────────────────────────────────────────
        out_path = self._output_dir / f"{output_name}.pptx"
        prs.save(str(out_path))

        slide_count = len(prs.slides)
        logger.info("PPT saved: %s (%d slides)", out_path, slide_count)

        await self.emit("ppt.created", {"path": str(out_path), "slides": slide_count})

        return TaskResult(
            task_id=context.task_id, agent_id=self.agent_id,
            outcome=ExecutionOutcome.SUCCESS,
            data={
                "file_path": str(out_path),
                "slide_count": slide_count,
                "title": title,
                "theme": theme_name,
                "topics": [t.get("title") for t in topics],
                "structure": [
                    "Title", "Table of Contents", "Introduction",
                    *[t.get("title", "Topic") for t in topics],
                    "Architecture Overview", "Summary", "Q & A",
                ],
            },
        )

    async def _export_presentation(self, context: TaskContext) -> TaskResult:
        file_path = context.parameters.get("file_path", "")
        path = Path(file_path)
        if not path.exists():
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=[f"File not found: {file_path}"],
            )
        size = path.stat().st_size
        return TaskResult(
            task_id=context.task_id, agent_id=self.agent_id,
            outcome=ExecutionOutcome.SUCCESS,
            data={"file_path": str(path), "size_bytes": size, "ready": True},
        )

    async def assess_self(self) -> dict[str, Any]:
        files = list(self._output_dir.glob("*.pptx"))
        return {
            "output_dir": str(self._output_dir),
            "presentations_created": len(files),
            "available_themes": list(PALETTE.keys()),
        }
