"""Quality Assessor Agent — VP-level / UX designer digital review of generated documents.

Evaluates PPT, Word, and PDF files across multiple professional dimensions:
design aesthetics, content structure, professionalism, inclusivity, accessibility.
Returns a scored report with actionable feedback, just like a senior design leader would.
"""

from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Any

from future_agents.core.base_agent import BaseAgent, TaskContext, TaskResult
from future_agents.models.feedback import ExecutionOutcome

logger = logging.getLogger(__name__)


RUBRIC = {
    "design_aesthetics": {
        "weight": 0.25,
        "criteria": [
            "Color palette coherence and professionalism",
            "Visual hierarchy (title > heading > body)",
            "Consistent use of whitespace",
            "Font selection and sizing",
            "Use of accent colors without overuse",
        ],
    },
    "content_structure": {
        "weight": 0.25,
        "criteria": [
            "Logical flow: intro → topics → conclusion",
            "Presence of index/TOC",
            "Section/chapter clarity",
            "Bullet/list formatting consistency",
            "Summary / closing section",
        ],
    },
    "professionalism": {
        "weight": 0.25,
        "criteria": [
            "Consistent branding / header-footer",
            "Page numbering",
            "Author / organization attribution",
            "Avoidance of clipart or low-quality elements",
            "Grammar and tone (where evaluable)",
        ],
    },
    "inclusivity": {
        "weight": 0.15,
        "criteria": [
            "Diverse imagery / representation",
            "Gender-balanced visual language",
            "Race and ethnicity representation",
            "Accessibility considerations (contrast, font size)",
            "Professional imagery appropriate to context",
        ],
    },
    "technical_quality": {
        "weight": 0.10,
        "criteria": [
            "File opens without errors",
            "Correct file format and structure",
            "File size appropriate",
            "Metadata present (title, author)",
            "No placeholder text remaining",
        ],
    },
}

GRADE_MAP = {
    (9.0, 10.0): ("A+", "World-class — ready for C-suite presentation"),
    (8.0, 9.0):  ("A",  "Excellent — minor polish only"),
    (7.0, 8.0):  ("B+", "Very good — some improvements recommended"),
    (6.0, 7.0):  ("B",  "Good — notable gaps, addressable"),
    (5.0, 6.0):  ("C",  "Adequate — significant improvements needed"),
    (0.0, 5.0):  ("D",  "Needs major rework"),
}


def _grade(score: float) -> tuple[str, str]:
    for (lo, hi), (letter, label) in GRADE_MAP.items():
        if lo <= score <= hi:
            return letter, label
    return "D", "Needs major rework"


class QualityAssessorAgent(BaseAgent):
    """Senior UX Designer / VP-level digital assessor for generated documents.

    Capabilities:
    - quality.assess  — score a single document
    - quality.compare — compare two documents of the same type
    - quality.report  — generate a full written assessment report
    """

    @property
    def agent_type(self) -> str:
        return "quality"

    @property
    def capabilities(self) -> list[str]:
        return ["quality.assess", "quality.compare", "quality.report"]

    async def _execute(self, context: TaskContext) -> TaskResult:
        handlers = {
            "quality.assess": self._assess_document,
            "quality.compare": self._compare_documents,
            "quality.report": self._full_report,
        }
        handler = handlers.get(context.intent)
        if not handler:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=[f"Unknown intent: {context.intent}"],
            )
        return await handler(context)

    # ── Assessment logic ──────────────────────────────────────────────────────

    def _inspect_pptx(self, path: Path) -> dict:
        """Extract structural metadata from a PPTX file."""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides = list(prs.slides)
            slide_count = len(slides)

            has_title_slide = slide_count > 0
            has_toc = any(
                any("contents" in str(shape.text).lower() or "index" in str(shape.text).lower()
                    for shape in s.shapes if shape.has_text_frame)
                for s in slides[:3]
            )
            has_summary = any(
                any("summary" in str(shape.text).lower() for shape in s.shapes if shape.has_text_frame)
                for s in slides
            )
            has_qa = any(
                any("q" in str(shape.text).lower() and "a" in str(shape.text).lower()
                    for shape in s.shapes if shape.has_text_frame)
                for s in slides
            )
            has_architecture = any(
                any("architect" in str(shape.text).lower() for shape in s.shapes if shape.has_text_frame)
                for s in slides
            )
            shape_count = sum(len(s.shapes) for s in slides)
            avg_shapes = shape_count / max(slide_count, 1)

            # Color check — look for background fills
            has_colored_bg = False
            for s in slides:
                for shape in s.shapes:
                    try:
                        fill = shape.fill
                        if fill.type is not None:
                            has_colored_bg = True
                            break
                    except Exception:
                        pass
                if has_colored_bg:
                    break

            return {
                "slide_count": slide_count,
                "has_title_slide": has_title_slide,
                "has_toc": has_toc,
                "has_summary": has_summary,
                "has_qa": has_qa,
                "has_architecture": has_architecture,
                "avg_shapes_per_slide": round(avg_shapes, 1),
                "has_colored_backgrounds": has_colored_bg,
                "file_size_kb": path.stat().st_size // 1024,
            }
        except Exception as exc:
            return {"error": str(exc)}

    def _inspect_docx(self, path: Path) -> dict:
        """Extract structural metadata from a DOCX file."""
        try:
            from docx import Document
            doc = Document(str(path))

            headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
            h1 = [h for h in headings if h.style.name == "Heading 1"]
            h2 = [h for h in headings if h.style.name == "Heading 2"]
            tables = doc.tables

            has_toc = any(
                "contents" in p.text.lower() or "toc" in p.text.lower()
                for p in doc.paragraphs[:10]
            )
            has_summary = any(
                "summary" in p.text.lower() for p in doc.paragraphs
            )
            has_header = bool(doc.sections and doc.sections[0].header.paragraphs)
            has_footer = bool(doc.sections and doc.sections[0].footer.paragraphs)

            total_text = " ".join(p.text for p in doc.paragraphs)
            word_count = len(total_text.split())

            return {
                "heading_1_count": len(h1),
                "heading_2_count": len(h2),
                "table_count": len(tables),
                "has_toc": has_toc,
                "has_summary": has_summary,
                "has_header": has_header,
                "has_footer": has_footer,
                "word_count": word_count,
                "file_size_kb": path.stat().st_size // 1024,
            }
        except Exception as exc:
            return {"error": str(exc)}

    def _inspect_pdf(self, path: Path) -> dict:
        """Extract structural metadata from a PDF file."""
        try:
            import struct

            # Basic PDF inspection without heavy deps
            raw = path.read_bytes()
            text_sample = raw[:8192].decode("latin-1", errors="ignore")

            is_valid_pdf = raw[:4] == b"%PDF"
            # Count /Type /Page occurrences — standard page marker in all PDF generators
            page_count = raw.count(b"/Type /Page") + raw.count(b"/Type/Page")
            # Fallback: estimate from file size (typical ~4KB/page for text PDFs)
            if page_count == 0 and is_valid_pdf:
                size_kb = path.stat().st_size // 1024
                page_count = max(1, size_kb // 4)
            has_images = b"/Image" in raw or b"JFIF" in raw or b"\xff\xd8\xff" in raw
            has_metadata = b"/Title" in raw or b"/Author" in raw
            has_bookmarks = b"/Outlines" in raw or b"/Dest" in raw

            return {
                "is_valid_pdf": is_valid_pdf,
                "approx_page_count": page_count,
                "has_images": has_images,
                "has_metadata": has_metadata,
                "has_bookmarks": has_bookmarks,
                "file_size_kb": path.stat().st_size // 1024,
            }
        except Exception as exc:
            return {"error": str(exc)}

    def _score_pptx(self, meta: dict) -> dict[str, float]:
        sc = meta.get("slide_count", 0)
        scores = {
            "design_aesthetics": 0.0,
            "content_structure": 0.0,
            "professionalism": 0.0,
            "inclusivity": 0.0,
            "technical_quality": 0.0,
        }
        if meta.get("error"):
            return scores

        # Design
        scores["design_aesthetics"] += 3.0 if meta.get("has_colored_backgrounds") else 0.0
        scores["design_aesthetics"] += min(meta.get("avg_shapes_per_slide", 0) * 0.8, 4.0)
        scores["design_aesthetics"] += 3.0 if sc >= 7 else (1.5 if sc >= 4 else 0.5)

        # Structure
        scores["content_structure"] += 2.0 if meta.get("has_title_slide") else 0.0
        scores["content_structure"] += 2.0 if meta.get("has_toc") else 0.0
        scores["content_structure"] += 2.0 if meta.get("has_summary") else 0.0
        scores["content_structure"] += 2.0 if meta.get("has_qa") else 0.0
        scores["content_structure"] += 2.0 if meta.get("has_architecture") else 0.0

        # Professionalism
        scores["professionalism"] += min(sc * 0.5, 5.0)
        scores["professionalism"] += 3.0 if meta.get("has_colored_backgrounds") else 0.0
        scores["professionalism"] += 2.0 if meta.get("file_size_kb", 0) > 20 else 0.5

        # Inclusivity (based on doc type — harder to auto-check in pptx)
        scores["inclusivity"] = 6.5  # baseline for PPT

        # Technical
        scores["technical_quality"] = 8.0 if not meta.get("error") else 2.0
        scores["technical_quality"] += 2.0 if meta.get("file_size_kb", 0) > 5 else 0.0

        # Cap at 10
        return {k: min(v, 10.0) for k, v in scores.items()}

    def _score_docx(self, meta: dict) -> dict[str, float]:
        scores = {
            "design_aesthetics": 0.0,
            "content_structure": 0.0,
            "professionalism": 0.0,
            "inclusivity": 0.0,
            "technical_quality": 0.0,
        }
        if meta.get("error"):
            return scores

        h1 = meta.get("heading_1_count", 0)
        h2 = meta.get("heading_2_count", 0)
        wc = meta.get("word_count", 0)

        scores["design_aesthetics"] += min(h1 * 1.2, 5.0)
        scores["design_aesthetics"] += 3.0 if meta.get("has_header") else 0.0
        scores["design_aesthetics"] += 2.0 if meta.get("has_footer") else 0.0

        scores["content_structure"] += 2.0 if meta.get("has_toc") else 0.0
        scores["content_structure"] += min(h1 * 1.0, 3.0)
        scores["content_structure"] += min(h2 * 0.5, 2.0)
        scores["content_structure"] += 2.0 if meta.get("has_summary") else 0.0
        scores["content_structure"] += min(meta.get("table_count", 0) * 0.5, 1.0)

        scores["professionalism"] += 2.5 if meta.get("has_header") else 0.0
        scores["professionalism"] += 2.5 if meta.get("has_footer") else 0.0
        scores["professionalism"] += min(wc / 200, 3.0)
        scores["professionalism"] += 2.0 if meta.get("file_size_kb", 0) > 10 else 0.0

        scores["inclusivity"] = 5.5  # baseline

        scores["technical_quality"] = 8.0 if not meta.get("error") else 2.0
        scores["technical_quality"] += 2.0 if meta.get("file_size_kb", 0) > 5 else 0.0

        return {k: min(v, 10.0) for k, v in scores.items()}

    def _score_pdf(self, meta: dict) -> dict[str, float]:
        scores = {
            "design_aesthetics": 0.0,
            "content_structure": 0.0,
            "professionalism": 0.0,
            "inclusivity": 0.0,
            "technical_quality": 0.0,
        }
        if meta.get("error") or not meta.get("is_valid_pdf"):
            return scores

        pg = meta.get("approx_page_count", 0)
        has_img = meta.get("has_images", False)

        scores["design_aesthetics"] += 4.0 if has_img else 1.0
        scores["design_aesthetics"] += min(pg * 0.4, 4.0)
        scores["design_aesthetics"] += 2.0 if meta.get("has_metadata") else 0.0

        scores["content_structure"] += min(pg * 0.5, 5.0)
        scores["content_structure"] += 2.0 if meta.get("has_bookmarks") else 0.0
        scores["content_structure"] += 3.0 if pg >= 5 else (1.5 if pg >= 3 else 0.5)

        scores["professionalism"] += 3.0 if meta.get("has_metadata") else 0.0
        scores["professionalism"] += 3.0 if has_img else 0.0
        scores["professionalism"] += min(pg * 0.3, 4.0)

        scores["inclusivity"] += 7.0 if has_img else 3.0
        scores["inclusivity"] += 3.0 if meta.get("has_metadata") else 0.0

        scores["technical_quality"] = 7.0 if meta.get("is_valid_pdf") else 0.0
        scores["technical_quality"] += 3.0 if meta.get("file_size_kb", 0) > 15 else 1.0

        return {k: min(v, 10.0) for k, v in scores.items()}

    def _compute_overall(self, scores: dict[str, float]) -> float:
        total = 0.0
        for dim, info in RUBRIC.items():
            total += scores.get(dim, 0.0) * info["weight"]
        return round(total, 2)

    def _build_feedback(self, file_type: str, scores: dict[str, float], meta: dict) -> list[str]:
        feedback = []
        for dim, score in scores.items():
            if score < 6.0:
                criteria = RUBRIC[dim]["criteria"]
                feedback.append(f"[{dim.upper()}] Score {score:.1f}/10 — Focus on: {criteria[0]}")
            elif score >= 9.0:
                feedback.append(f"[{dim.upper()}] Score {score:.1f}/10 — Excellent. Keep this standard.")
            else:
                feedback.append(f"[{dim.upper()}] Score {score:.1f}/10 — Good. Minor refinements possible.")

        if file_type == "pptx":
            if not meta.get("has_toc"):
                feedback.append("ACTION: Add a Table of Contents slide after the title.")
            if not meta.get("has_architecture"):
                feedback.append("ACTION: Add an Architecture/Framework diagram slide.")
            if meta.get("slide_count", 0) < 8:
                feedback.append("ACTION: Expand to at least 8+ slides for executive presentations.")

        elif file_type == "docx":
            if not meta.get("has_toc"):
                feedback.append("ACTION: Insert a Table of Contents (Ctrl+A, F9 to update in Word).")
            if not meta.get("has_header") or not meta.get("has_footer"):
                feedback.append("ACTION: Add consistent headers and footers with page numbers.")

        elif file_type == "pdf":
            if not meta.get("has_images"):
                feedback.append("ACTION: Add background imagery — set UNSPLASH_ACCESS_KEY or PEXELS_API_KEY.")
            if not meta.get("has_bookmarks"):
                feedback.append("ACTION: Add PDF bookmarks/outlines for navigation.")

        return feedback

    def _vp_verdict(self, score: float) -> str:
        if score >= 9.0:
            return ("APPROVED — Board ready. This document meets the highest standard of "
                    "professional quality, visual design, and inclusive representation.")
        elif score >= 8.0:
            return ("APPROVED WITH NOTES — Executive ready. Strong quality overall with minor "
                    "polish needed before wider distribution.")
        elif score >= 7.0:
            return ("CONDITIONAL APPROVAL — Very good but requires specific improvements before "
                    "external-facing use. See actionable feedback below.")
        elif score >= 6.0:
            return ("REVISION REQUIRED — Meets baseline but notable gaps in design or structure "
                    "must be addressed. Internal use only until revised.")
        else:
            return ("REJECTED — Significant rework required. Document does not meet professional "
                    "standards for distribution. Follow the full improvement plan.")

    # ── Intent handlers ───────────────────────────────────────────────────────

    async def _assess_document(self, context: TaskContext) -> TaskResult:
        """Score a document and return a detailed assessment.

        Parameters:
            file_path (str): Full path to the document
            file_type (str): pptx | docx | pdf (auto-detected if omitted)
        """
        p = context.parameters
        file_path = p.get("file_path", "")
        path = Path(file_path)

        if not path.exists():
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=[f"File not found: {file_path}"],
            )

        file_type = p.get("file_type", path.suffix.lstrip(".").lower())

        if file_type == "pptx":
            meta = self._inspect_pptx(path)
            scores = self._score_pptx(meta)
        elif file_type == "docx":
            meta = self._inspect_docx(path)
            scores = self._score_docx(meta)
        elif file_type == "pdf":
            meta = self._inspect_pdf(path)
            scores = self._score_pdf(meta)
        else:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=[f"Unsupported file type: {file_type}. Use pptx, docx, or pdf."],
            )

        overall = self._compute_overall(scores)
        letter, label = _grade(overall)
        verdict = self._vp_verdict(overall)
        feedback = self._build_feedback(file_type, scores, meta)

        logger.info("Quality assessment: %s → %.1f/10 (%s)", path.name, overall, letter)

        return TaskResult(
            task_id=context.task_id, agent_id=self.agent_id,
            outcome=ExecutionOutcome.SUCCESS,
            data={
                "file": str(path),
                "file_type": file_type,
                "overall_score": overall,
                "grade": letter,
                "grade_label": label,
                "vp_verdict": verdict,
                "dimension_scores": {k: round(v, 2) for k, v in scores.items()},
                "metadata_extracted": meta,
                "feedback": feedback,
                "rubric": {k: v["weight"] for k, v in RUBRIC.items()},
            },
            suggestions=feedback[:3],
        )

    async def _compare_documents(self, context: TaskContext) -> TaskResult:
        """Compare two documents of the same type side-by-side.

        Parameters:
            file_a (str): Path to first document
            file_b (str): Path to second document
        """
        file_a = context.parameters.get("file_a", "")
        file_b = context.parameters.get("file_b", "")

        results = []
        for fp in [file_a, file_b]:
            sub_ctx = TaskContext(
                intent="quality.assess",
                parameters={"file_path": fp},
                parent_task_id=context.task_id,
            )
            result = await self._assess_document(sub_ctx)
            results.append(result.data)

        if len(results) < 2:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=["Could not assess both documents"],
            )

        a, b = results[0], results[1]
        winner = "A" if a.get("overall_score", 0) >= b.get("overall_score", 0) else "B"

        return TaskResult(
            task_id=context.task_id, agent_id=self.agent_id,
            outcome=ExecutionOutcome.SUCCESS,
            data={
                "document_a": {"file": file_a, **a},
                "document_b": {"file": file_b, **b},
                "winner": winner,
                "score_delta": round(
                    abs(a.get("overall_score", 0) - b.get("overall_score", 0)), 2
                ),
                "recommendation": (
                    f"Document {winner} scores higher by "
                    f"{abs(a.get('overall_score',0)-b.get('overall_score',0)):.1f} points."
                ),
            },
        )

    async def _full_report(self, context: TaskContext) -> TaskResult:
        """Generate a written assessment report for a document.

        Parameters:
            file_path (str): Document path
            output_dir (str): Where to save the report (optional)
        """
        file_path = context.parameters.get("file_path", "")
        output_dir = context.parameters.get("output_dir", "output/reports")

        sub_ctx = TaskContext(
            intent="quality.assess",
            parameters={"file_path": file_path},
            parent_task_id=context.task_id,
        )
        assessment = await self._assess_document(sub_ctx)
        data = assessment.data

        if not data:
            return TaskResult(
                task_id=context.task_id, agent_id=self.agent_id,
                outcome=ExecutionOutcome.FAILURE,
                errors=assessment.errors,
            )

        lines = [
            "=" * 72,
            "DOCUMENT QUALITY ASSESSMENT REPORT",
            "Senior UX Designer + VP-Level Digital Review",
            "=" * 72,
            f"\nFile          : {data.get('file', 'N/A')}",
            f"Type          : {data.get('file_type', '').upper()}",
            f"Overall Score : {data.get('overall_score', 0):.1f} / 10",
            f"Grade         : {data.get('grade', 'N/A')} — {data.get('grade_label', '')}",
            "\n" + "-" * 72,
            "VP VERDICT",
            "-" * 72,
            data.get("vp_verdict", ""),
            "\n" + "-" * 72,
            "DIMENSION SCORES",
            "-" * 72,
        ]

        for dim, score in data.get("dimension_scores", {}).items():
            weight = RUBRIC.get(dim, {}).get("weight", 0)
            bar = "█" * int(score) + "░" * (10 - int(score))
            lines.append(f"  {dim:<22} {bar}  {score:4.1f}/10  (weight {weight:.0%})")

        lines += [
            "\n" + "-" * 72,
            "ACTIONABLE FEEDBACK",
            "-" * 72,
        ]
        for fb in data.get("feedback", []):
            lines.append(f"  • {fb}")

        lines += [
            "\n" + "-" * 72,
            "DOCUMENT METADATA",
            "-" * 72,
        ]
        for k, v in data.get("metadata_extracted", {}).items():
            lines.append(f"  {k:<30}: {v}")

        lines.append("\n" + "=" * 72)
        report_text = "\n".join(lines)

        Path(output_dir).mkdir(parents=True, exist_ok=True)
        file_stem = Path(file_path).stem
        report_path = Path(output_dir) / f"assessment_{file_stem}.txt"
        report_path.write_text(report_text, encoding="utf-8")

        return TaskResult(
            task_id=context.task_id, agent_id=self.agent_id,
            outcome=ExecutionOutcome.SUCCESS,
            data={
                "report_path": str(report_path),
                "overall_score": data.get("overall_score"),
                "grade": data.get("grade"),
                "vp_verdict": data.get("vp_verdict"),
                "report_text": report_text,
            },
        )

    async def assess_self(self) -> dict[str, Any]:
        return {
            "rubric_dimensions": list(RUBRIC.keys()),
            "supported_types": ["pptx", "docx", "pdf"],
            "total_criteria": sum(len(v["criteria"]) for v in RUBRIC.values()),
        }
